import os
from convnet_drawer import *
from pptx import Presentation
from pptx.shapes.connector import Connector
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat


def get_or_add_ln(self):
    return self._element.spPr.get_or_add_ln()


Connector.get_or_add_ln = get_or_add_ln


class MyPresentation:
    def __init__(self):
        self.presentation = Presentation(os.path.join(os.path.dirname(__file__), "template.pptx"))
        self.slide_layout = self.presentation.slide_layouts[6]
        self.slide = self.presentation.slides.add_slide(self.slide_layout)
        self.shapes = self.slide.shapes

    def add_line(self, x1, y1, x2, y2, color, width, dasharray):
        connector = self.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))

        if not hasattr(connector, "ln"):
            connector.ln = connector.get_or_add_ln()

        line = LineFormat(connector)
        line.width = Pt(width)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*color)

        if dasharray == 1:
            line.dash_style = MSO_LINE.SQUARE_DOT
        elif dasharray == 2:
            line.dash_style = MSO_LINE.DASH

    def add_text(self, x, y, body, color, size):
        # TODO: set color
        textbox = self.shapes.add_textbox(Pt(x), Pt(y), Pt(0), Pt(0))
        textbox.text = body
        text_frame = textbox.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = text_frame.paragraphs[0]
        font = p.font
        font.name = 'arial'
        font.size = Pt(size)
        p.alignment = PP_ALIGN.CENTER

    def save_pptx(self, filename):
        self.presentation.save(filename)


def save_model_to_pptx(model, filename, shift=0.0, scale=1.0):
    model.build()
    presentation = MyPresentation()

    for feature_map in model.feature_maps + model.layers:
        for obj in feature_map.objects:
            if isinstance(obj, Line):
                presentation.add_line(obj.x1 * scale, obj.y1 * scale + shift, obj.x2 * scale, obj.y2 * scale + shift,
                                      obj.color, obj.width, obj.dasharray)
            elif isinstance(obj, Text):
                presentation.add_text(obj.x * scale, obj.y * scale + shift, obj.body, obj.color, obj.size * scale)

    presentation.save_pptx(filename)
